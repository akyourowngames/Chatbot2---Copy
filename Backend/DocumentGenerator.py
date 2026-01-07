"""
Advanced Document Generator - Beast Mode
=========================================
Generate professional, content-rich PDFs and PowerPoint presentations with AI integration.
Features:
- AI-generated content for Reports, Articles, and Presentations
- Premium styling with cover pages and headers/footers
- Support regarding images, tables, and diverse layouts
"""

from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle, Frame, PageTemplate
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY, TA_RIGHT
from reportlab.pdfgen import canvas

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import os
import json
import re
from datetime import datetime
from typing import List, Dict, Any, Optional

try:
    from Backend.LLM import ChatCompletion
except ImportError:
    # Fallback for testing independently
    def ChatCompletion(messages, model, text_only): return "AI Content Generation Unavailable"

class DocumentGenerator:
    def __init__(self):
        self.output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "Data", "Documents")
        os.makedirs(self.output_dir, exist_ok=True)
        
        # Define brand colors
        self.primary_color = colors.HexColor('#2563EB')     # Royal Blue
        self.secondary_color = colors.HexColor('#1E40AF')   # Darker Blue
        self.accent_color = colors.HexColor('#F59E0B')      # Amber
        self.text_color = colors.HexColor('#1F2937')        # Gray-800
        self.light_bg = colors.HexColor('#F3F4F6')          # Gray-100

    def _get_llm_content(self, topic: str, doc_type: str, detailed: bool = True) -> Dict[str, Any]:
        """
        Generate structured content using the LLM.
        """
        print(f"Generating content for {doc_type} on '{topic}'...")
        
        prompt = f"""
        You are a professional content creator. Create detailed content for a {doc_type} about "{topic}".
        Return the response strictly as a valid JSON object.
        
        Structure for 'report' or 'article':
        {{
            "title": "Impactful Title",
            "subtitle": "Engaging Subtitle",
            "author": "JARVIS AI",
            "sections": [
                {{
                    "heading": "Section Heading",
                    "content": "Detailed paragraph content...",
                    "type": "paragraph" | "bullet" | "table", 
                    "data": ["point 1", "point 2"] (if bullet) or [["Header 1", "Header 2"], ["Row 1 Col 1", "Row 1 Col 2"]] (if table)
                }}
            ]
        }}

        Structure for 'presentation':
        {{
            "title": "Presentation Title",
            "slides": [
                {{
                    "title": "Slide Title",
                    "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
                    "layout": "bullet" | "title" | "two_column"
                }}
            ]
        }}
        
        Make the content informative, professional, and comprehensive. {'Include formatting like bolding or lists where appropriate.' if detailed else ''}
        """
        
        try:
            response = ChatCompletion(
                messages=[{"role": "system", "content": "You are a JSON-generating AI. Output ONLY JSON."}, {"role": "user", "content": prompt}],
                model="llama-3.3-70b-versatile",
                text_only=True
            )
            
            # extract json block if needed
            if "```json" in response:
                response = response.split("```json")[1].split("```")[0]
            elif "```" in response:
                 response = response.split("```")[1].split("```")[0]
                 
            return json.loads(response.strip())
        except Exception as e:
            print(f"Error generating AI content: {e}")
            # Fallback content
            return {
                "title": f"{doc_type.title()} on {topic}",
                "sections": [{"heading": "Error", "content": "Could not generate content.", "type": "paragraph"}]
            }

    # ==================== PDF GENERATION ====================

    def _add_header_footer(self, canvas, doc):
        """Add styled header and footer to PDF pages"""
        canvas.saveState()
        
        # Header line
        canvas.setStrokeColor(self.primary_color)
        canvas.setLineWidth(2)
        canvas.line(72, letter[1] - 50, letter[0] - 72, letter[1] - 50)
        
        # Footer
        page_num = canvas.getPageNumber()
        date_str = datetime.now().strftime("%B %d, %Y")
        
        canvas.setFont('Helvetica', 9)
        canvas.setFillColor(colors.grey)
        
        # Left footer (Date)
        canvas.drawString(72, 30, date_str)
        
        # Center footer (Brand)
        canvas.drawCentredString(letter[0]/2, 30, "Generated by JARVIS AI")
        
        # Right footer (Page Number)
        canvas.drawRightString(letter[0] - 72, 30, f"Page {page_num}")
        
        canvas.restoreState()

    def generate_pdf(self, topic: str, content: Dict[str, Any] = None, filename: str = None, user_id: str = None) -> Dict[str, str]:
        """
        Generate a beast-level PDF document.
        If content is None, it generates it via LLM.
        
        Args:
            topic: Document topic
            content: Pre-generated content (optional)
            filename: Custom filename (optional)
            user_id: User ID for user-specific storage (optional)
        
        Returns:
            Dict with 'filepath' (local) and 'url' (cloud or local path)
        """
        if content is None:
            content = self._get_llm_content(topic, "report")
            
        title = content.get("title", topic)
        if not filename:
            filename = f"{title.replace(' ', '_')[:50]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = os.path.join(self.output_dir, filename)

        doc = SimpleDocTemplate(filepath, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
        elements = []
        styles = getSampleStyleSheet()

        # --- Custom Styles ---
        title_style = ParagraphStyle('MainTitle', parent=styles['Title'], fontSize=32, spaceAfter=20, textColor=self.primary_color, fontName='Helvetica-Bold')
        subtitle_style = ParagraphStyle('SubTitle', parent=styles['Heading2'], fontSize=16, alignment=TA_CENTER, spaceAfter=40, textColor=colors.grey)
        
        heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'], fontSize=18, spaceBefore=20, spaceAfter=10, textColor=self.secondary_color, fontName='Helvetica-Bold')
        body_style = ParagraphStyle('CustomBody', parent=styles['Normal'], fontSize=11, leading=16, spaceAfter=10, alignment=TA_JUSTIFY)
        bullet_style = ParagraphStyle('CustomBullet', parent=body_style, leftIndent=20, bulletIndent=10)

        # --- Cover Page ---
        elements.append(Spacer(1, 2*inch))
        elements.append(Paragraph(title, title_style))
        if "subtitle" in content:
            elements.append(Paragraph(content["subtitle"], subtitle_style))
        
        elements.append(Spacer(1, 1*inch))
        elements.append(Paragraph(f"Created for: {os.environ.get('USERNAME', 'User')}", ParagraphStyle('Meta', parent=body_style, alignment=TA_CENTER)))
        elements.append(Paragraph(f"Date: {datetime.now().strftime('%B %d, %Y')}", ParagraphStyle('Meta', parent=body_style, alignment=TA_CENTER)))
        elements.append(PageBreak())

        # --- Content Pages ---
        sections = content.get("sections", [])
        for section in sections:
            heading = section.get("heading")
            text = section.get("content", "")
            sec_type = section.get("type", "paragraph")
            data = section.get("data", [])

            if heading:
                elements.append(Paragraph(heading, heading_style))

            if sec_type == "paragraph":
                # Handle basic Markdown bolding (**text**)
                text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)
                elements.append(Paragraph(text, body_style))
            
            elif sec_type == "bullet":
                if text: elements.append(Paragraph(text, body_style))
                for item in data:
                    item = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', item)
                    elements.append(Paragraph(f"• {item}", bullet_style))
            
            elif sec_type == "table":
                if text: elements.append(Paragraph(text, body_style))
                if data:
                    t = Table(data)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), self.primary_color),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), self.light_bg),
                        ('GRID', (0,0), (-1,-1), 1, colors.white),
                        ('PADDING', (0,0), (-1,-1), 8),
                    ]))
                    elements.append(t)
            
            elements.append(Spacer(1, 12))

        doc.build(elements, onFirstPage=self._add_header_footer, onLaterPages=self._add_header_footer)
        
        # Always upload to Supabase Storage (kai-images bucket)
        pdf_url = f"/data/Documents/{filename}"  # Default fallback
        try:
            from Backend.SupabaseDB import supabase_db
            if supabase_db:
                print(f"[DocumentGenerator] Uploading PDF to Supabase (kai-images bucket)...")
                # Pass user_id for user-specific storage
                cloud_url = supabase_db.upload_pdf(filepath, folder='documents', user_id=user_id)
                if cloud_url:
                    print(f"[DocumentGenerator] Uploaded to: {cloud_url}")
                    pdf_url = cloud_url
                    # Keep local file as backup (don't delete)
                else:
                    print(f"[DocumentGenerator] Upload failed, using local path")
        except Exception as e:
            print(f"[DocumentGenerator] Supabase upload error: {e}, using local path")
        
        return {
            "filepath": filepath,
            "url": pdf_url,
            "filename": filename,
            "title": title
        }

    # ==================== PPTX GENERATION ====================

    def generate_presentation(self, topic: str, content: Dict[str, Any] = None, filename: str = None, user_id: str = None) -> str:
        """
        Generate a beast-level PowerPoint presentation.
        """
        if content is None:
            content = self._get_llm_content(topic, "presentation")

        title = content.get("title", topic)
        if not filename:
            filename = f"{title.replace(' ', '_')[:50]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(self.output_dir, filename)

        prs = Presentation()
        
        # --- Title Slide ---
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"Presented by JARVIS AI\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Style Title Slide
        title_shape.text_frame.paragraphs[0].font.size = Pt(54)
        title_shape.text_frame.paragraphs[0].font.name = 'Calibri'
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235) # Brand Blue

        # --- Content Slides ---
        for i, slide_data in enumerate(content.get("slides", [])):
            layout_idx = 1 # Title and Content
            if slide_data.get("layout") == "title": layout_idx = 0
            elif slide_data.get("layout") == "two_column": layout_idx = 3
            
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            
            # Set Title
            title_shape = slide.shapes.title
            title_shape.text = slide_data.get("title", f"Slide {i+1}")
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175) # Dark Blue
            
            # Set Content
            slide_content = slide_data.get("content", [])
            if layout_idx == 1: # Standard Bullet
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for point in slide_content:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.space_after = Pt(10)

        prs.save(filepath)

        # Upload to Supabase if available
        try:
            from Backend.SupabaseDB import supabase_db
            if supabase_db:
                cloud_url = supabase_db.upload_file(filepath, f"documents/{filename}" if not user_id else f"{user_id}/documents/{filename}", bucket='kai-images', content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                print(f"[DOC] Uploaded PPTX to Supabase: {cloud_url}")
        except Exception as e:
            print(f"[WARN] Failed to upload PPTX to Supabase: {e}")

        return filepath

    def create_document(self, topic: str, doc_type: str = "pdf", user_id: str = None) -> str:
        """
        Unified entry point to create a document or presentation from a topic.
        """
        if "presentation" in doc_type.lower() or "ppt" in doc_type.lower():
            return self.generate_presentation(topic, user_id=user_id)
        else:
            return self.generate_pdf(topic, user_id=user_id)

# Global Instance
document_generator = DocumentGenerator()

if __name__ == "__main__":
    # Test Run
    print("Testing PDF creation...")
    path = document_generator.create_document("The Future of Artificial Intelligence", "pdf")
    print(f"Created: {path}")
    
    print("\nTesting PPT creation...")
    path_ppt = document_generator.create_document("Strategies for Growth 2025", "ppt")
    print(f"Created: {path_ppt}")
